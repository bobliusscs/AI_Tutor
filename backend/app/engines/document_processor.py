"""
文档处理工具 - 支持PDF、图片、Word、PPT等文档的解析和转换
"""
import os
import base64
import re
from typing import List, Tuple, Optional, Dict, Any
from pathlib import Path
import io


def clean_text_for_utf8(text: str) -> str:
    """
    清理文本中的 surrogate 字符，确保可以安全编码为 UTF-8
    
    PDF 提取的文本可能包含数学符号等字符的 incomplete surrogate pair，
    这些字符在 JSON 序列化/UTF-8 编码时会报错。此函数移除或替换这些字符。
    
    Args:
        text: 原始文本
        
    Returns:
        清理后的文本
    """
    if not text:
        return text
    
    # 方法1：使用 surrogateescape 处理，然后替换为 �
    # 这会保留所有可编码字符，将 surrogate 替换为替换字符
    try:
        return text.encode('utf-8', errors='surrogateescape').decode('utf-8', errors='replace')
    except Exception:
        # 兜底：直接移除 surrogate 范围的字符 (U+D800-U+DFFF)
        return re.sub(r'[\ud800-\udfff]', '', text)

# PDF处理
try:
    import PyPDF2
    import pdfplumber
    HAS_PDF_SUPPORT = True
except ImportError:
    HAS_PDF_SUPPORT = False
    print("警告: PyPDF2 或 pdfplumber 未安装，PDF处理功能将受限")

# Word处理
try:
    from docx import Document as DocxDocument
    HAS_DOCX_SUPPORT = True
except ImportError:
    HAS_DOCX_SUPPORT = False
    print("警告: python-docx 未安装，Word文档处理功能将受限")

# PPT处理
try:
    from pptx import Presentation
    HAS_PPTX_SUPPORT = True
except ImportError:
    HAS_PPTX_SUPPORT = False
    print("警告: python-pptx 未安装，PPT处理功能将受限")

# 图片处理
try:
    from PIL import Image
    HAS_PIL_SUPPORT = True
except ImportError:
    HAS_PIL_SUPPORT = False
    print("警告: Pillow 未安装，图片处理功能将受限")


class DocumentProcessor:
    """文档处理器"""
    
    def __init__(self, upload_dir: str = None):
        """
        初始化文档处理器

        Args:
            upload_dir: 上传文件目录路径（默认与main.py保持一致）
        """
        if upload_dir:
            self.upload_dir = upload_dir
        else:
            # 默认路径：与main.py保持一致
            self.upload_dir = os.path.join(
                os.path.dirname(os.path.dirname(os.path.dirname(__file__))),
                "app", "uploads", "materials"
            )
    
    def process_file(self, file_path: str) -> Dict[str, Any]:
        """
        处理任意支持的文档文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            处理结果字典，包含:
            - type: 文档类型 (scanned_pdf/text_pdf/word/powerpoint/image/text)
            - text: 提取的文字内容（如果有）
            - images: 图片路径列表（如果是扫描版PDF或图片）
            - is_scanned: 是否是扫描版PDF
            - page_count: 页数
        """
        file_path = self._get_absolute_path(file_path)
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.pdf':
            return self.process_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return self.process_word(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self.process_powerpoint(file_path)
        elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.webp']:
            return self.process_image(file_path)
        elif ext == '.txt':
            return self.process_text(file_path)
        else:
            raise ValueError(f"不支持的文件类型: {ext}")
    
    def process_pdf(self, file_path: str) -> Dict[str, Any]:
        """
        处理PDF文件，区分扫描版和文字版
        
        Args:
            file_path: PDF文件路径
            
        Returns:
            处理结果字典
        """
        if not HAS_PDF_SUPPORT:
            raise ImportError("PDF处理库未安装")
        
        file_path = self._get_absolute_path(file_path)
        
        result = {
            "type": "pdf",
            "text": "",
            "images": [],
            "is_scanned": False,
            "page_count": 0,
            "file_path": file_path
        }
        
        print(f"[PDF处理] 开始处理: {file_path}")
        
        # 首先尝试提取文字，判断是否为扫描版
        extracted_text = self._extract_pdf_text(file_path)
        print(f"[PDF处理] 提取到文字长度: {len(extracted_text)}")
        
        # 判断是否为扫描版PDF（文字内容很少或为空）
        is_scanned = self._is_scanned_pdf(extracted_text, file_path)
        
        if is_scanned:
            result["is_scanned"] = True
            result["type"] = "scanned_pdf"
            print(f"[PDF处理] 判定为扫描版PDF，开始转换为图片...")
            # 扫描版PDF需要转换为图片
            result["images"] = self._convert_pdf_to_images(file_path)
            print(f"[PDF处理] 图片转换完成，共 {len(result['images'])} 张图片")
        else:
            result["text"] = extracted_text
            print(f"[PDF处理] 判定为文字版PDF")
        
        # 获取页数
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            result["page_count"] = len(pdf_reader.pages)
        
        print(f"[PDF处理] 处理完成，结果: type={result['type']}, images={len(result['images'])}, text_len={len(result['text'])}")
        return result
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """
        从PDF提取文字
        
        Args:
            file_path: PDF文件路径
            
        Returns:
            提取的文字内容（已清理 surrogate 字符）
        """
        if not HAS_PDF_SUPPORT:
            return ""
        
        text_parts = []
        
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
        except Exception as e:
            print(f"pdfplumber提取失败: {e}，尝试使用PyPDF2")
            try:
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    for page in pdf_reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(page_text)
            except Exception as e2:
                print(f"PyPDF2提取也失败: {e2}")
        
        raw_text = "\n\n".join(text_parts)
        # 清理 surrogate 字符，确保可以安全编码为 UTF-8
        return clean_text_for_utf8(raw_text)
    
    def _is_scanned_pdf(self, text: str, file_path: str, min_text_length: int = 100) -> bool:
        """
        判断PDF是否为扫描版

        判断依据：
        1. 提取的文字内容非常少（每页平均字符数过低）
        2. 总字符数过少
        3. 可读字符比例过低（大量乱码）
        4. 文字分布异常（如每页文字数量差异极大）
        5. 页面大小与文字量不匹配

        Args:
            text: 提取的文字
            file_path: 文件路径
            min_text_length: 最小文字长度阈值

        Returns:
            是否为扫描版
        """
        # 清理文字
        clean_text = text.strip()

        print(f"[PDF检测] 开始分析文件: {os.path.basename(file_path)}")

        # 统计信息收集
        page_analysis = []
        try:
            import PyPDF2
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                page_count = len(pdf_reader.pages)

                # 逐页分析文字分布
                with pdfplumber.open(file_path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        page_text = page.extract_text() or ""
                        # 统计该页的有效字符数（中英文、数字、标点）
                        import re
                        valid_chars = len(re.findall(r'[\u4e00-\u9fff\u0041-\u005a\u0061-\u007a\u0030-\u0039\u4e00-\u9fa5，。、！？；：""''（）【】《》\d]', page_text))
                        page_analysis.append({
                            "page": i + 1,
                            "valid_chars": valid_chars,
                            "total_chars": len(page_text)
                        })

                print(f"[PDF检测] 文件: {os.path.basename(file_path)}, 页数: {page_count}, "
                      f"提取文字长度: {len(clean_text)}, 每页平均: {len(clean_text) / max(page_count, 1):.1f}字符")

                # 1. 每页平均字符数检查
                chars_per_page = len(clean_text) / max(page_count, 1)
                if chars_per_page < 30:  # 降低阈值，更严格判断
                    print(f"[PDF检测] 判定为扫描版PDF（每页平均字符数过低: {chars_per_page:.1f} < 30）")
                    return True

                # 2. 检查页面间文字分布的差异性
                if len(page_analysis) > 1:
                    valid_counts = [p["valid_chars"] for p in page_analysis]
                    max_chars = max(valid_counts) if valid_counts else 0
                    min_chars = min(valid_counts) if valid_counts else 0

                    # 如果最大差异超过50倍，可能是扫描版（有些页有文字，有些页是纯图片）
                    if max_chars > 0 and max_chars / max_chars > 50:
                        print(f"[PDF检测] 判定为扫描版PDF（页面文字分布差异过大: {min_chars}-{max_chars}）")
                        return True

                    # 检查是否有大量页面文字极少（少于20个有效字符）
                    sparse_pages = sum(1 for c in valid_counts if c < 20)
                    if sparse_pages / len(valid_counts) > 0.7:  # 超过70%的页面文字极少
                        print(f"[PDF检测] 判定为扫描版PDF（{sparse_pages}/{len(valid_counts)}页文字极少）")
                        return True

        except Exception as e:
            print(f"[PDF检测] 获取页数失败: {e}")

        # 3. 如果提取的文字太少，判定为扫描版
        if len(clean_text) < min_text_length:
            print(f"[PDF检测] 判定为扫描版PDF（总字符数过低: {len(clean_text)} < {min_text_length}）")
            return True

        # 4. 计算可读字符比例
        import re
        # 统计可读中英文字符和常见标点
        readable_chars = len(re.findall(r'[\u4e00-\u9fff\u0041-\u005a\u0061-\u007a\u0030-\u0039，。、！？；：""''（）【】《》\s]', text))
        total_chars = len(text)

        if total_chars > 0:
            readable_ratio = readable_chars / total_chars
            print(f"[PDF检测] 可读字符比例: {readable_ratio:.2%} ({readable_chars}/{total_chars})")
            # 如果可读字符比例低于20%，判定为扫描版（更严格）
            if readable_ratio < 0.2:
                print(f"[PDF检测] 判定为扫描版PDF（可读字符比例过低: {readable_ratio:.2%} < 20%）")
                return True

        # 5. 检查是否包含明显的乱码模式
        # 统计非标准空白字符和控制字符
        non_standard_chars = len(re.findall(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', text))
        if total_chars > 0 and non_standard_chars / total_chars > 0.05:  # 超过5%非标准字符
            print(f"[PDF检测] 判定为扫描版PDF（非标准字符过多: {non_standard_chars}/{total_chars}）")
            return True

        # 6. 检查句子完整性
        # 文字版PDF通常包含完整句子，扫描版可能只有碎片化的文字
        sentences = re.findall(r'[。！？；\n]+', text)
        if len(sentences) > 0:
            # 计算平均每句话的长度
            avg_sentence_len = len(text) / len(sentences)
            if avg_sentence_len < 5:  # 平均每句话太短，可能是乱码
                print(f"[PDF检测] 判定为扫描版PDF（句子过于破碎: 平均{avg_sentence_len:.1f}字符/句）")
                return True

        print(f"[PDF检测] 判定为文字版PDF")
        return False
    
    def _convert_pdf_to_images(self, file_path: str, dpi: int = 150) -> List[str]:
        """
        将PDF转换为图片
        
        Args:
            file_path: PDF文件路径
            dpi: 转换分辨率
            
        Returns:
            图片文件路径列表
        """
        print(f"[图片转换] 开始转换PDF: {file_path}, DPI: {dpi}")
        
        # 检查库是否可用
        has_pymupdf = False
        has_pdf2image = False
        
        try:
            import fitz  # PyMuPDF
            has_pymupdf = True
            print(f"[图片转换] PyMuPDF 可用")
        except ImportError:
            print(f"[图片转换] PyMuPDF 不可用")
        
        try:
            from pdf2image import convert_from_path
            import poppler
            has_pdf2image = True
            print(f"[图片转换] pdf2image + poppler 可用")
        except ImportError:
            print(f"[图片转换] pdf2image 或 poppler 不可用")
        
        if not has_pymupdf and not has_pdf2image:
            print(f"[图片转换] 错误: PyMuPDF 和 pdf2image 都未安装，无法将PDF转换为图片")
            print(f"[图片转换] 请安装: pip install pymupdf 或 pip install pdf2image poppler")
            return []
        
        images_dir = os.path.join(self.upload_dir, "_temp_images")
        os.makedirs(images_dir, exist_ok=True)
        
        image_paths = []
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        
        # 尝试使用 PyMuPDF
        if has_pymupdf:
            try:
                print(f"[图片转换] 使用 PyMuPDF 转换...")
                doc = fitz.open(file_path)
                page_count = len(doc)
                print(f"[图片转换] PDF总页数: {page_count}")
                
                for page_num in range(page_count):
                    page = doc[page_num]
                    # 渲染页面为图片
                    mat = fitz.Matrix(dpi / 72, dpi / 72)
                    pix = page.get_pixmap(matrix=mat)
                    
                    # 保存图片
                    img_path = os.path.join(images_dir, f"{base_name}_page_{page_num + 1}.png")
                    pix.save(img_path)
                    image_paths.append(img_path)
                    print(f"[图片转换] 已转换第 {page_num + 1}/{page_count} 页 -> {img_path}")
                
                doc.close()
                print(f"[图片转换] PyMuPDF 转换完成，共 {len(image_paths)} 张图片")
                return image_paths
            except Exception as e:
                print(f"[图片转换] PyMuPDF转换失败: {e}")
        
        # 尝试使用 pdf2image
        if has_pdf2image:
            try:
                print(f"[图片转换] 尝试使用 pdf2image 转换...")
                images = convert_from_path(file_path, dpi=dpi)
                print(f"[图片转换] pdf2image 转换完成，共 {len(images)} 张图片")
                
                for i, img in enumerate(images):
                    img_path = os.path.join(images_dir, f"{base_name}_page_{i + 1}.png")
                    img.save(img_path, "PNG")
                    image_paths.append(img_path)
                
                print(f"[图片转换] pdf2image 转换完成，共 {len(image_paths)} 张图片")
                return image_paths
            except Exception as e2:
                print(f"[图片转换] pdf2image转换也失败: {e2}")
        
        print(f"[图片转换] 所有转换方法都失败，返回空列表")
        return image_paths
    
    def process_word(self, file_path: str) -> Dict[str, Any]:
        """
        处理Word文档
        
        Args:
            file_path: Word文件路径
            
        Returns:
            处理结果字典
        """
        if not HAS_DOCX_SUPPORT:
            raise ImportError("python-docx 未安装")
        
        file_path = self._get_absolute_path(file_path)
        
        result = {
            "type": "word",
            "text": "",
            "images": [],
            "page_count": 1,
            "file_path": file_path
        }
        
        doc = DocxDocument(file_path)
        
        # 提取所有段落的文字
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        # 提取表格中的文字
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text)
        
        result["text"] = clean_text_for_utf8("\n\n".join(paragraphs))
        
        # 统计段落数作为页数估计
        result["page_count"] = max(1, len(paragraphs) // 30)
        
        return result
    
    def process_powerpoint(self, file_path: str) -> Dict[str, Any]:
        """
        处理PowerPoint演示文稿
        
        Args:
            file_path: PPT文件路径
            
        Returns:
            处理结果字典
        """
        if not HAS_PPTX_SUPPORT:
            raise ImportError("python-pptx 未安装")
        
        file_path = self._get_absolute_path(file_path)
        
        result = {
            "type": "powerpoint",
            "text": "",
            "images": [],
            "page_count": 0,
            "file_path": file_path
        }
        
        prs = Presentation(file_path)
        result["page_count"] = len(prs.slides)
        
        # 提取每张幻灯片的文字
        slide_texts = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        slide_text.append(shape.text)
            
            if slide_text:
                slide_texts.append(f"=== 第{i+1}页 ===\n" + "\n".join(slide_text))
        
        result["text"] = clean_text_for_utf8("\n\n".join(slide_texts))
        
        return result
    
    def process_image(self, file_path: str) -> Dict[str, Any]:
        """
        处理图片文件
        
        Args:
            file_path: 图片文件路径
            
        Returns:
            处理结果字典
        """
        if not HAS_PIL_SUPPORT:
            raise ImportError("Pillow 未安装")
        
        file_path = self._get_absolute_path(file_path)
        
        result = {
            "type": "image",
            "text": "",  # 图片本身不做OCR，交给多模态模型处理
            "images": [file_path],
            "page_count": 1,
            "file_path": file_path
        }
        
        return result
    
    def process_text(self, file_path: str) -> Dict[str, Any]:
        """
        处理纯文本文件
        
        Args:
            file_path: 文本文件路径
            
        Returns:
            处理结果字典
        """
        file_path = self._get_absolute_path(file_path)
        
        result = {
            "type": "text",
            "text": "",
            "images": [],
            "page_count": 1,
            "file_path": file_path
        }
        
        with open(file_path, 'r', encoding='utf-8') as f:
            result["text"] = clean_text_for_utf8(f.read())
        
        return result
    
    def _get_absolute_path(self, file_path: str) -> str:
        """
        获取文件的绝对路径

        Args:
            file_path: 文件路径（可能是相对路径或URL路径）

        Returns:
            绝对路径
        """
        # 保存原始文件名
        original_filename = os.path.basename(file_path)
        
        # 处理URL路径 (/uploads/...) - 在Windows上这也被视为绝对路径，需要先处理
        if file_path.startswith('/uploads/'):
            file_path = file_path.replace('/uploads/', '')
        
        # 如果已经是绝对路径（Windows驱动器路径如 C:\... 或 Linux /home/...）
        if os.path.isabs(file_path):
            return file_path
        
        # 获取项目根目录
        # backend/app/engines/document_processor.py
        # -> backend/app/engines/
        # -> backend/app/
        # -> backend/
        # -> 项目根目录
        engine_dir = os.path.dirname(os.path.dirname(__file__))  # app/
        app_dir = os.path.dirname(engine_dir)  # backend/
        project_root = os.path.dirname(app_dir)  # AI_Tutor/
        
        # 尝试多个可能的基础目录（包含所有历史可能的路径）
        possible_dirs = [
            # 正确的配置位置
            os.path.join(app_dir, "app", "uploads", "materials"),  # backend/app/uploads/materials
            # 可能的错误位置（旧配置）
            os.path.join(app_dir, "app", "api", "uploads", "materials"),  # backend/app/api/uploads/materials
            os.path.join(app_dir, "api", "uploads", "materials"),  # backend/api/uploads/materials（兼容旧路径）
            os.path.join(app_dir, "uploads", "materials"),  # backend/uploads/materials
            os.path.join(app_dir, "app", "uploads"),  # backend/app/uploads
            os.path.join(app_dir, "app", "api", "uploads"),  # backend/app/api/uploads
            os.path.join(app_dir, "api", "uploads"),  # backend/api/uploads
            os.path.join(app_dir, "uploads"),  # backend/uploads
            # 项目根目录下的
            os.path.join(project_root, "backend", "uploads", "materials"),
            os.path.join(project_root, "uploads", "materials"),
            # 当前工作目录
            os.getcwd(),
        ]
        
        # 去重并过滤不存在的目录
        seen = set()
        clean_dirs = []
        for d in possible_dirs:
            if d not in seen and os.path.exists(os.path.dirname(d) if not os.path.isdir(d) else d):
                seen.add(d)
                if os.path.isdir(d):
                    clean_dirs.append(d)
        
        # 尝试找到文件 - 在所有可能的目录中搜索同名文件
        for base_dir in clean_dirs:
            # 直接在目录中搜索同名文件
            if os.path.isdir(base_dir):
                for root, dirs, files in os.walk(base_dir):
                    if original_filename in files:
                        full_path = os.path.join(root, original_filename)
                        print(f"[路径解析] 在目录 {base_dir} 中找到文件: {full_path}")
                        return full_path
        
        # 如果搜索都没找到，尝试基于配置的目录
        result = os.path.join(self.upload_dir, file_path)
        print(f"[路径解析] 文件 '{original_filename}' 未找到，返回计算的路径: {result}")
        return result
    
    def image_to_base64(self, image_path: str) -> str:
        """
        将图片转换为Base64编码
        
        Args:
            image_path: 图片路径
            
        Returns:
            Base64编码字符串
        """
        image_path = self._get_absolute_path(image_path)
        
        with open(image_path, 'rb') as f:
            return base64.b64encode(f.read()).decode('utf-8')
    
    def split_text(self, text: str, chunk_size: int = 2000, overlap: int = 200) -> List[str]:
        """
        将长文本分割为多个小块
        
        Args:
            text: 原始文本
            chunk_size: 每块的目标字符数
            overlap: 块之间的重叠字符数
            
        Returns:
            分割后的文本块列表
        """
        if len(text) <= chunk_size:
            return [text] if text.strip() else []
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # 尽量在句子或段落边界分割
            if end < len(text):
                # 向前查找最近的换行符或句号
                search_range = text[max(0, end - 200):min(len(text), end + 200)]
                newline_pos = search_range.rfind('\n')
                period_pos = search_range.rfind('。')
                
                if newline_pos != -1:
                    end = max(start + chunk_size, end - 200 + newline_pos)
                elif period_pos != -1:
                    end = max(start + chunk_size, end - 200 + period_pos + 1)
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - overlap
            if start <= chunks[-1].__len__() if chunks else 0:
                start = end
        
        return [c for c in chunks if c.strip()]
    
    def cleanup_temp_images(self, image_paths: List[str]):
        """
        清理临时生成的图片文件
        
        Args:
            image_paths: 图片路径列表
        """
        for path in image_paths:
            try:
                if os.path.exists(path):
                    os.remove(path)
            except Exception as e:
                print(f"删除临时图片失败: {path}, 错误: {e}")


# 全局处理器实例
processor = DocumentProcessor()
